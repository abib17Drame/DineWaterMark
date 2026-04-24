"""
Traitement PPTX : extraction images → suppression watermark → reconstruction.

Pipeline :
1. Réception du fichier PPTX
2. Extraction des images de chaque slide
3. Conversion en PDF lossless via img2pdf (optionnel)
4. Suppression du watermark (inpainting ou colonne par colonne)
5. Reconstruction en PPTX propre via PyMuPDF + python-pptx
"""
import os
import io
import logging
import time
import shutil
from pathlib import Path

import fitz  # PyMuPDF
import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu

from app.services.watermark.suppresseur import supprimer_watermark

logger = logging.getLogger(__name__)


def traiter_pptx(
    chemin_entree: str,
    chemin_sortie: str,
    mode_debug: bool = False,
    callback_progression: callable = None,
    mode_watermark: str = "auto",
) -> dict:
    """
    Traite un fichier PPTX : supprime le watermark de chaque slide.

    Stratégie :
    - Chaque slide d'un PPTX NotebookLM contient typiquement
      une seule image de fond (la slide rendue).
    - On extrait cette image, on supprime le watermark, et on la remplace.

    Args:
        chemin_entree: Chemin vers le PPTX d'entrée
        chemin_sortie: Chemin vers le PPTX de sortie
        mode_debug: Activer les infos de debug
        callback_progression: Fonction(slide_courante, total_slides)

    Returns:
        Dictionnaire de statistiques
    """
    debut = time.time()

    presentation = Presentation(chemin_entree)
    nombre_slides = len(presentation.slides)

    logger.info(
        "Début du traitement PPTX : %s (%d slides)",
        os.path.basename(chemin_entree), nombre_slides,
    )

    slides_avec_watermark = 0
    resultats_slides = []

    for index_slide, slide in enumerate(presentation.slides):
        infos_slide = {"slide": index_slide + 1, "watermark_detecte": False}
        preview_orig = None
        preview_nette = None

        # Chercher les images dans la slide
        for forme in slide.shapes:
            if not forme.shape_type == 13:  # 13 = MSO_SHAPE_TYPE.PICTURE
                continue

            # Extraire les données de l'image
            blob_image = forme.image.blob
            type_contenu = forme.image.content_type

            # Convertir en array NumPy via PIL
            image_pil = Image.open(io.BytesIO(blob_image))
            image_np = np.array(image_pil)

            # Convertir en BGR pour OpenCV
            if len(image_np.shape) == 2:
                # Image en niveaux de gris
                image_bgr = cv2.cvtColor(image_np, cv2.COLOR_GRAY2BGR)
            elif image_np.shape[2] == 4:
                # Image RGBA
                image_bgr = cv2.cvtColor(image_np, cv2.COLOR_RGBA2BGR)
            else:
                # Image RGB
                image_bgr = cv2.cvtColor(image_np, cv2.COLOR_RGB2BGR)

            # Supprimer le watermark
            image_nettoyee, infos_traitement = supprimer_watermark(
                image_bgr, mode_debug=mode_debug, mode_watermark=mode_watermark
            )

            preview_orig = image_bgr
            preview_nette = image_nettoyee

            if infos_traitement.get("watermark_detecte"):
                slides_avec_watermark += 1
                infos_slide["watermark_detecte"] = True
                infos_slide.update(infos_traitement)

                # Reconvertir en RGB et sauvegarder temporairement pour l'intégration
                image_nettoyee_rgb = cv2.cvtColor(image_nettoyee, cv2.COLOR_BGR2RGB)
                img_pil_nettoyee = Image.fromarray(image_nettoyee_rgb)

                chemin_temp_img = os.path.join(
                    os.path.dirname(chemin_sortie),
                    f"_temp_clean_{index_slide}_{os.getpid()}.png"
                )
                img_pil_nettoyee.save(chemin_temp_img, format="PNG")

                # Récupérer les dimensions et positions
                gauche = forme.left
                haut = forme.top
                largeur = forme.width
                hauteur = forme.height

                # Supprimer l'ancienne forme de l'arbre XML pour alléger le fichier
                sp = forme._element
                sp.getparent().remove(sp)

                # Insérer la nouvelle image propre
                slide.shapes.add_picture(
                    chemin_temp_img, gauche, haut, largeur, hauteur
                )

                # Nettoyage
                if os.path.exists(chemin_temp_img):
                    os.remove(chemin_temp_img)

        resultats_slides.append(infos_slide)

        if callback_progression:
            try:
                import base64

                if preview_orig is not None and preview_nette is not None:
                    h, w = preview_orig.shape[:2]
                    ratio = 300 / w
                    dim = (300, int(h * ratio))

                    small_orig = cv2.resize(preview_orig, dim, interpolation=cv2.INTER_AREA)
                    _, buf_orig = cv2.imencode('.jpg', small_orig, [cv2.IMWRITE_JPEG_QUALITY, 50])
                    b64_orig = base64.b64encode(buf_orig).decode('utf-8')

                    small_nette = cv2.resize(preview_nette, dim, interpolation=cv2.INTER_AREA)
                    _, buf_nette = cv2.imencode('.jpg', small_nette, [cv2.IMWRITE_JPEG_QUALITY, 50])
                    b64_nette = base64.b64encode(buf_nette).decode('utf-8')

                    callback_progression(index_slide + 1, nombre_slides, b64_orig, b64_nette)
                else:
                    callback_progression(index_slide + 1, nombre_slides)
            except Exception as e:
                logger.error(f"Erreur preview base64 PPTX: {e}")
                callback_progression(index_slide + 1, nombre_slides)

        logger.info(
            "Slide %d/%d traitée (watermark=%s)",
            index_slide + 1, nombre_slides,
            "OUI" if infos_slide["watermark_detecte"] else "NON",
        )

    # Sauvegarder le PPTX modifié
    presentation.save(chemin_sortie)

    duree_ms = int((time.time() - debut) * 1000)

    resultat_global = {
        "status": "succes",
        "type_fichier": "pptx",
        "nombre_slides": nombre_slides,
        "slides_avec_watermark": slides_avec_watermark,
        "temps_traitement_ms": duree_ms,
        "details_slides": resultats_slides if mode_debug else None,
    }

    logger.info(
        "PPTX traité avec succès en %dms — %d/%d slides avec watermark",
        duree_ms, slides_avec_watermark, nombre_slides,
    )

    return resultat_global


def fusionner_pptx(
    chemins_entree: list[str],
    chemin_sortie: str,
) -> dict:
    """
    Fusionne plusieurs fichiers PPTX en un seul (ordre préservé).

    Utile pour NotebookLM qui limite à ~15 slides par fichier.

    Args:
        chemins_entree: Liste des chemins PPTX à fusionner
        chemin_sortie: Chemin du PPTX fusionné

    Returns:
        Statistiques de fusion
    """
    debut = time.time()

    if not chemins_entree:
        raise ValueError("Aucun fichier PPTX à fusionner")

    # Utiliser le premier fichier comme base
    presentation_base = Presentation(chemins_entree[0])
    nombre_total_slides = len(presentation_base.slides)

    for chemin_supplementaire in chemins_entree[1:]:
        pres_supplementaire = Presentation(chemin_supplementaire)

        for slide_source in pres_supplementaire.slides:
            # Copier la slide dans la présentation de base
            layout_slide = presentation_base.slide_layouts[6]  # Layout vide
            nouvelle_slide = presentation_base.slides.add_slide(layout_slide)

            # Copier toutes les formes de la slide source
            for forme in slide_source.shapes:
                if forme.shape_type == 13:  # Image
                    blob = forme.image.blob
                    gauche = forme.left
                    haut = forme.top
                    largeur_forme = forme.width
                    hauteur_forme = forme.height

                    # Sauvegarder le blob en fichier temporaire
                    chemin_temp_img = os.path.join(
                        os.path.dirname(chemin_sortie),
                        f"_temp_merge_{os.getpid()}.png"
                    )
                    with open(chemin_temp_img, "wb") as f:
                        f.write(blob)

                    nouvelle_slide.shapes.add_picture(
                        chemin_temp_img, gauche, haut,
                        largeur_forme, hauteur_forme
                    )

                    # Nettoyage du fichier temporaire
                    if os.path.exists(chemin_temp_img):
                        os.remove(chemin_temp_img)

            nombre_total_slides += 1

    presentation_base.save(chemin_sortie)

    duree_ms = int((time.time() - debut) * 1000)

    logger.info(
        "Fusion PPTX terminée : %d fichiers → %d slides en %dms",
        len(chemins_entree), nombre_total_slides, duree_ms,
    )

    return {
        "status": "succes",
        "fichiers_fusionnes": len(chemins_entree),
        "nombre_total_slides": nombre_total_slides,
        "temps_traitement_ms": duree_ms,
    }
 